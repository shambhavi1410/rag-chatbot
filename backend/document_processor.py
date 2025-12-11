import os
import asyncio
from typing import List
from langchain.text_splitter import RecursiveCharacterTextSplitter

# NO langchain_community imports (prevents pwd error)
import pytesseract
from PIL import Image
import docx
import pptx
import io
import PyPDF2

# -------------------------------------------------------
# Configure Tesseract path for Windows
# -------------------------------------------------------
try:
    import platform
    if platform.system() == 'Windows':
        tesseract_paths = [
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
        ]
        for path in tesseract_paths:
            if os.path.exists(path):
                pytesseract.pytesseract.tesseract_cmd = path
                break
except:
    pass


class DocumentProcessor:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
    
    async def process_document(self, file_path: str, file_extension: str) -> List[str]:
        """Process different document types and extract text asynchronously"""
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            None, self._process_document_sync, file_path, file_extension
        )
    
    def _process_document_sync(self, file_path: str, file_extension: str) -> List[str]:
        """Synchronous document processing"""
        texts = []
        
        try:
            if file_extension == ".pdf":
                texts = self._process_pdf(file_path)
            elif file_extension in [".pptx", ".ppt"]:
                texts = self._process_ppt(file_path)
            elif file_extension in [".docx", ".doc"]:
                texts = self._process_docx(file_path)
            elif file_extension.lower() in [".jpg", ".jpeg", ".png", ".gif", ".bmp"]:
                texts = self._process_image(file_path)
            elif file_extension == ".txt":
                texts = self._process_txt(file_path)
            else:
                texts = self._process_txt(file_path)

            if texts:
                chunks = self.text_splitter.split_text("\n".join(texts))
                return chunks
            else:
                return []
        except Exception as e:
            print(f"Error processing document: {e}")
            return []
    
    # -------------------------------------------------------
    # PDF PROCESSING WITHOUT LANGCHAIN
    # -------------------------------------------------------
    def _process_pdf(self, file_path: str) -> List[str]:
        texts = []
        try:
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    texts.append(page.extract_text() or "")
        except Exception as e:
            print(f"PDF Error: {e}")
        return texts

    # -------------------------------------------------------
    # POWERPOINT PROCESSING WITHOUT LANGCHAIN
    # -------------------------------------------------------
    def _process_ppt(self, file_path: str) -> List[str]:
        texts = []
        try:
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                if slide_text:
                    texts.append("\n".join(slide_text))
        except Exception as e:
            print(f"PPT Error: {e}")
        return texts
    
    # -------------------------------------------------------
    # WORD DOCUMENT PROCESSING WITHOUT LANGCHAIN
    # -------------------------------------------------------
    def _process_docx(self, file_path: str) -> List[str]:
        texts = []
        try:
            doc = docx.Document(file_path)
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    texts.append(paragraph.text)
        except Exception as e:
            print(f"DOCX Error: {e}")
        return texts
    
    # -------------------------------------------------------
    # IMAGE OCR PROCESSING
    # -------------------------------------------------------
    def _process_image(self, file_path: str) -> List[str]:
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image, lang='eng+hin')  
            return [text] if text.strip() else []
        except Exception as e:
            print(f"OCR Error: {e}")
            return []
    
    # -------------------------------------------------------
    # TEXT FILE PROCESSING
    # -------------------------------------------------------
    def _process_txt(self, file_path: str) -> List[str]:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return [f.read()]
        except Exception as e:
            print(f"TXT Error: {e}")
            return []

